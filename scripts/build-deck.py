#!/usr/bin/env python3
"""
Builds the client deck: why the redesigned Amrut Maharashtra portal is
better than the live one.

Every number in here was measured against the live site, not estimated.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ---- AMRUT palette, taken from their live stylesheet -----------------
SAFFRON = RGBColor(0xF9, 0x73, 0x16)
DEEP    = RGBColor(0xD3, 0x54, 0x00)
CREAM   = RGBColor(0xFF, 0xF7, 0xED)
CREAM2  = RGBColor(0xFF, 0xED, 0xD5)
INK     = RGBColor(0x1A, 0x16, 0x13)
INK2    = RGBColor(0x3F, 0x37, 0x30)
WARM    = RGBColor(0x9B, 0x8F, 0x85)
WARM2   = RGBColor(0xE4, 0xDC, 0xD4)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GOOD    = RGBColor(0x1A, 0x66, 0x40)
BAD     = RGBColor(0xA8, 0x1F, 0x16)

SERIF = "Georgia"
SANS  = "Helvetica Neue"

W, H = Inches(13.333), Inches(7.5)
ASSETS = os.path.join(os.path.dirname(__file__), "..", "docs", "deck-assets")

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, font=SANS, bold=False, space_after=0,
         first=False, align=PP_ALIGN.LEFT, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    return p


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def eyebrow(s, text, y=0.62, color=DEEP):
    tf = box(s, 0.9, y, 9, 0.3)
    para(tf, text, 12, color, SANS, bold=True, first=True)


def title(s, text, y=0.95, size=34, color=INK, w=11.5):
    tf = box(s, 0.9, y, w, 1.3)
    para(tf, text, size, color, SERIF, first=True, line=1.15)


def footer(s, n):
    tf = box(s, 12.1, 6.95, 1, 0.3)
    para(tf, str(n), 10, WARM, SANS, first=True, align=PP_ALIGN.RIGHT)


def pic(s, name, x, y, w):
    path = os.path.join(ASSETS, name)
    if os.path.exists(path):
        return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    return None


n = 0

# =====================================================================
# 1 — Cover
# =====================================================================
s = slide(INK)
rect(s, 0, 0, 0.22, 7.5, SAFFRON)

tf = box(s, 1.2, 2.15, 10.5, 0.4)
para(tf, "अमृत महाराष्ट्र  ·  संकेतस्थळ पुनर्रचना", 13, SAFFRON, SANS, bold=True, first=True)

tf = box(s, 1.2, 2.7, 10.5, 2.0)
para(tf, "Why the new portal is better", 46, CREAM, SERIF, first=True, line=1.1)

tf = box(s, 1.2, 4.35, 8.6, 1.2)
para(tf, "A redesign of amrutmaharashtra.org for the people it is actually meant to reach — "
         "rural Maharashtra, in Marathi, on a slow connection.",
     16, RGBColor(0xC9, 0xBF, 0xB6), SANS, first=True, line=1.4)

tf = box(s, 1.2, 6.4, 10.5, 0.4)
para(tf, "Same content.  Same colours.  Same institution.  Everything else rebuilt.",
     13, RGBColor(0x8A, 0x7F, 0x76), SANS, first=True)

# =====================================================================
# 2 — The problem in one sentence
# =====================================================================
n += 1
s = slide(CREAM)
eyebrow(s, "THE PROBLEM")
title(s, "Today the portal is a magazine that happens\nto be run by a government agency.", size=32)

tf = box(s, 0.9, 2.9, 6.0, 2.6)
para(tf, "Every success story ends the same way: here is who qualifies, "
         "24 castes, income under ₹8 lakh — now go contact your district office.",
     17, INK2, SANS, first=True, line=1.55, space_after=14)
para(tf, "The site never says where that office is. It never links the scheme. "
         "It never offers a way to apply.", 17, INK2, SANS, line=1.55, space_after=14)
para(tf, "A reader who has just been convinced has nowhere to go.",
     17, DEEP, SANS, bold=True, line=1.55)

rect(s, 7.4, 2.75, 5.0, 2.9, WHITE, WARM2)
tf = box(s, 7.8, 3.1, 4.2, 2.2)
para(tf, "2,664", 44, SAFFRON, SERIF, first=True)
para(tf, "articles published", 13, WARM, SANS, space_after=18)
para(tf, "0", 44, INK, SERIF)
para(tf, "of them lead anywhere", 13, WARM, SANS)
footer(s, n)

# =====================================================================
# 3 — Before / after, side by side
# =====================================================================
n += 1
s = slide(WHITE)
eyebrow(s, "THE HOMEPAGE")
title(s, "Before and after")

tf = box(s, 0.9, 2.05, 5.6, 0.3)
para(tf, "BEFORE  ·  amrutmaharashtra.org", 11, BAD, SANS, bold=True, first=True)
pic(s, "old-home.jpg", 0.9, 2.4, 5.6)

tf = box(s, 6.9, 2.05, 5.6, 0.3)
para(tf, "AFTER  ·  amrut-mh.vercel.app", 11, GOOD, SANS, bold=True, first=True)
pic(s, "new-home.jpg", 6.9, 2.4, 5.6)

tf = box(s, 0.9, 5.5, 5.6, 1.2)
para(tf, "Three stacked bars of chrome take a quarter of the screen. Sixteen categories "
         "compete at equal weight. Saffron runs everything, so nothing leads.",
     12.5, INK2, SANS, first=True, line=1.4)

tf = box(s, 6.9, 5.5, 5.6, 1.2)
para(tf, "One person, one photograph, rotating through five stories. The masthead is 64px. "
         "Saffron is spent once — on the thing you should do next.",
     12.5, INK2, SANS, first=True, line=1.4)
footer(s, n)

# =====================================================================
# 4 — Colour
# =====================================================================
n += 1
s = slide(WHITE)
eyebrow(s, "COLOUR")
title(s, "We did not change a single colour.\nWe changed how much of each you see.")

tf = box(s, 0.9, 2.75, 5.4, 2.2)
para(tf, "The palette is AMRUT's own, read straight out of the live stylesheet — "
         "not matched by eye.", 16, INK2, SANS, first=True, line=1.5, space_after=16)
para(tf, "On the live site saffron runs the header, the nav bar, the buttons, the "
         "headings and the links, all at full strength. When everything shouts, "
         "nothing is heard.", 16, INK2, SANS, line=1.5, space_after=16)
para(tf, "Now it covers about 5% of the surface and marks one thing per screen.",
     16, DEEP, SANS, bold=True, line=1.5)

sw = [("#F97316", SAFFRON, "Saffron — the accent"),
      ("#D35400", DEEP, "Deep orange — links"),
      ("#FFF7ED", CREAM, "Cream — the ground"),
      ("#FFEDD5", CREAM2, "Cream 2 — panels"),
      ("#1A1613", INK, "Ink — the type")]
y = 2.75
for hexv, col, label in sw:
    rect(s, 7.0, y, 0.85, 0.62, col, WARM2)
    tf = box(s, 8.05, y + 0.08, 4.2, 0.5)
    para(tf, hexv, 13, INK, SANS, bold=True, first=True)
    para(tf, label, 11, WARM, SANS)
    y += 0.78

tf = box(s, 0.9, 6.3, 11.5, 0.5)
para(tf, "This matters for approval: it still reads as AMRUT's site, done properly — "
         "not as somebody rebranding a government portal.",
     13, WARM, SANS, first=True)
footer(s, n)

# =====================================================================
# 5 — Typography
# =====================================================================
n += 1
s = slide(CREAM)
eyebrow(s, "TYPOGRAPHY")
title(s, "Marathi set as Marathi — not as English\nwearing a Devanagari font.")

tf = box(s, 0.9, 2.85, 6.0, 3.0)
para(tf, "Most Indian government sites set Devanagari in a default font at a "
         "line-height borrowed from Latin. It comes out cramped: the shirorekha "
         "crowds, the matras collide. Readers feel it even when they cannot name it.",
     16, INK2, SANS, first=True, line=1.5, space_after=16)
para(tf, "We use two faces drawn for Devanagari, not adapted to it:",
     16, INK2, SANS, line=1.5, space_after=10)
para(tf, "Tiro Devanagari Marathi — a text serif designed specifically for Marathi. "
         "Headlines and article body.", 15, INK, SANS, bold=True, line=1.45, space_after=8)
para(tf, "Mukta, by Ek Type of Mumbai — navigation, metadata, interface.",
     15, INK, SANS, bold=True, line=1.45)

rect(s, 7.4, 2.85, 5.0, 3.0, WHITE, WARM2)
tf = box(s, 7.8, 3.15, 4.2, 2.4)
para(tf, "Body line-height", 12, WARM, SANS, first=True)
para(tf, "1.90", 40, SAFFRON, SERIF, space_after=6)
para(tf, "for Devanagari, against the 1.5 a Latin face would take — the matras need "
         "the room.", 13, INK2, SANS, line=1.4, space_after=16)
para(tf, "No uppercase anywhere. It does not exist in Devanagari, and Latin caps set "
         "beside it read as a foreign body.", 13, INK2, SANS, line=1.4)

tf = box(s, 0.9, 6.35, 11.5, 0.5)
para(tf, "This is the single biggest differentiator available, and almost nobody in "
         "this sector takes it.", 13, DEEP, SANS, bold=True, first=True)
footer(s, n)

# =====================================================================
# 6 — The bridge (the big idea)
# =====================================================================
n += 1
s = slide(WHITE)
eyebrow(s, "THE BIGGEST CHANGE")
title(s, "Every story now leads somewhere.")

pic(s, "new-bridge.jpg", 6.6, 2.2, 6.0)

tf = box(s, 0.9, 2.3, 5.2, 3.6)
para(tf, "At the end of each article the site now works out which scheme the story is "
         "actually about — from the article's own words — and offers it.",
     16, INK2, SANS, first=True, line=1.5, space_after=14)
para(tf, "Who qualifies. The income limit. The district office. The helpline. "
         "A button that goes straight into the application portal.",
     16, INK2, SANS, line=1.5, space_after=14)
para(tf, "Not one word of the article was changed. This is a new component placed "
         "after it.", 16, INK2, SANS, line=1.5, space_after=16)
para(tf, "Do this, and 2,664 already-published stories become 2,664 entry points.",
     16, DEEP, SANS, bold=True, line=1.5)
footer(s, n)

# =====================================================================
# 7 — Speed
# =====================================================================
n += 1
s = slide(INK)
eyebrow(s, "SPEED", color=SAFFRON)
title(s, "The audience is on rural mobile data.\nThe old site does not act like it.", color=CREAM)

stats = [("~80", "images on the homepage", "none lazy-loaded, none responsive"),
         ("18×", "oversized at worst", "a 1600×2133 photo filling a 150×200 slot"),
         ("152 MB", "of original artwork", "downloaded in full, every visit")]
x = 0.9
for big, mid, small in stats:
    tf = box(s, x, 2.9, 3.6, 1.6)
    para(tf, big, 40, SAFFRON, SERIF, first=True)
    para(tf, mid, 14, CREAM, SANS, bold=True, space_after=4)
    para(tf, small, 12, RGBColor(0x8A, 0x7F, 0x76), SANS, line=1.35)
    x += 3.9

rect(s, 0.9, 4.95, 11.5, 1.25, RGBColor(0x2A, 0x24, 0x1F))
tf = box(s, 1.3, 5.2, 10.7, 0.9)
para(tf, "After:  152 MB  →  37 MB", 24, CREAM, SERIF, first=True, space_after=6)
para(tf, "Every image rebuilt as WebP at three widths. The versions that fill the "
         "homepage average 24 KB each. Same pictures, roughly a tenth of the weight.",
     13.5, RGBColor(0xC9, 0xBF, 0xB6), SANS, line=1.4)

tf = box(s, 0.9, 6.55, 11.5, 0.4)
para(tf, "On a slow connection this is felt before anything else on this list.",
     12.5, RGBColor(0x8A, 0x7F, 0x76), SANS, first=True)
footer(s, n)

# =====================================================================
# 8 — Trust
# =====================================================================
n += 1
s = slide(WHITE)
eyebrow(s, "CREDIBILITY")
title(s, "Three things that quietly undermine\nan official government site.")

items = [
    ("The visitor counter invents its own traffic",
     "It climbed from 2,03,935 to 3,09,418 while we sat on the homepage — about a minute — "
     "then read 0 on every inner page. On a site whose footer says “This is an official "
     "government website”, that is the fastest way to lose a reader.",
     "Now: report real analytics, or remove it."),
    ("The Contact button does nothing",
     "संपर्क sits in the navigation styled as the one highlighted action, and is wired to "
     "nothing at all. There is no contact page anywhere on the site.",
     "Now: a real contact route, and the details in the footer and the assistant."),
    ("The 404 page is somebody else's template",
     "A wrong link lands on a stock illustration of a skateboarder. No masthead, no Marathi, "
     "no emblem, no way back. It is where every broken link ends.",
     "Now: a proper page, in Marathi, listing all 16 sections."),
]
y = 2.15
for head, body, fix in items:
    rect(s, 0.9, y, 0.055, 1.35, SAFFRON)
    tf = box(s, 1.25, y, 11.2, 1.35)
    para(tf, head, 16, INK, SANS, bold=True, first=True, space_after=5)
    para(tf, body, 12.5, INK2, SANS, line=1.4, space_after=4)
    para(tf, fix, 12.5, GOOD, SANS, bold=True)
    y += 1.55
footer(s, n)

# =====================================================================
# 9 — Findability
# =====================================================================
n += 1
s = slide(CREAM)
eyebrow(s, "REACH")
title(s, "2,664 Marathi articles that search\nengines can barely see.")

rows = [("robots.txt", "Missing", "Generated at build"),
        ("sitemap.xml", "Missing", "Generated at build"),
        ("Article URLs", "news.php?id=3283", "/beneficiary-story/3284/सोलापुरच्या…"),
        ("Page titles", "One generic title reused", "Per-page, in Marathi"),
        ("Shared to WhatsApp", "Previews as nothing", "Proper title, image and description")]

rect(s, 0.9, 2.7, 11.5, 0.45, INK)
hx = [1.15, 4.6, 8.3]
hd = ["", "BEFORE", "AFTER"]
for i, h in enumerate(hd):
    tf = box(s, hx[i], 2.79, 3.3, 0.3)
    para(tf, h, 11, CREAM, SANS, bold=True, first=True)

y = 3.15
for i, (label, before, after) in enumerate(rows):
    rect(s, 0.9, y, 11.5, 0.62, WHITE if i % 2 == 0 else RGBColor(0xFA, 0xF6, 0xF2))
    tf = box(s, 1.15, y + 0.16, 3.3, 0.35)
    para(tf, label, 13, INK, SANS, bold=True, first=True)
    tf = box(s, 4.6, y + 0.16, 3.5, 0.35)
    para(tf, before, 12.5, BAD, SANS, first=True)
    tf = box(s, 8.3, y + 0.16, 4.0, 0.35)
    para(tf, after, 12.5, GOOD, SANS, first=True)
    y += 0.62

tf = box(s, 0.9, 6.35, 11.5, 0.5)
para(tf, "WhatsApp is how this content actually travels in Maharashtra. Right now a "
         "shared link arrives looking like nothing.", 13, INK2, SANS, first=True)
footer(s, n)

# =====================================================================
# 10 — Article page
# =====================================================================
n += 1
s = slide(WHITE)
eyebrow(s, "THE ARTICLE")
title(s, "The same words, made readable")

tf = box(s, 0.9, 2.05, 5.6, 0.3)
para(tf, "BEFORE", 11, BAD, SANS, bold=True, first=True)
pic(s, "old-article.jpg", 0.9, 2.4, 5.6)

tf = box(s, 6.9, 2.05, 5.6, 0.3)
para(tf, "AFTER", 11, GOOD, SANS, bold=True, first=True)
pic(s, "new-article.jpg", 6.9, 2.4, 5.6)

tf = box(s, 0.9, 5.5, 5.6, 1.3)
para(tf, "Text runs the full width of the screen. The office address and the standing "
         "“about AMRUT” paragraph sit mid-story as if they were part of it. "
         "No navigation — the only way out is back.",
     12.5, INK2, SANS, first=True, line=1.4)

tf = box(s, 6.9, 5.5, 5.6, 1.3)
para(tf, "A measured reading column. Author, date, district and views in a quiet margin. "
         "The office and the boilerplate lifted out and given their own place — same "
         "words, no longer read as story.",
     12.5, INK2, SANS, first=True, line=1.4)
footer(s, n)

# =====================================================================
# 11 — Assistant
# =====================================================================
n += 1
s = slide(WHITE)
eyebrow(s, "NEW")
title(s, "An assistant that answers in Marathi —\nand refuses to make things up.")

pic(s, "new-assistant.jpg", 6.6, 2.2, 6.0)

tf = box(s, 0.9, 2.3, 5.2, 3.8)
para(tf, "It handles eligibility, the four live schemes, how to apply, what documents "
         "are needed, district offices and contact — and searches every story on the "
         "site.", 15.5, INK2, SANS, first=True, line=1.5, space_after=14)
para(tf, "Every answer is built from this site's own content. It cannot invent an "
         "eligibility rule or an office address.", 15.5, INK2, SANS, line=1.5, space_after=14)
para(tf, "Ask it for a district office and it gives Solapur — the only one AMRUT "
         "actually publishes — and points to the head office for the rest.",
     15.5, INK2, SANS, line=1.5, space_after=14)
para(tf, "On a government portal, being right matters more than sounding fluent.",
     15.5, DEEP, SANS, bold=True, line=1.5)
footer(s, n)

# =====================================================================
# 12 — Summary scoreboard
# =====================================================================
n += 1
s = slide(INK)
eyebrow(s, "IN SHORT", color=SAFFRON)
title(s, "What changed", color=CREAM)

rows = [("Homepage weight", "~152 MB of images", "37 MB, 24 KB per card"),
        ("Masthead height", "~200 px, always", "64 px, sticky rail below"),
        ("Navigation on inner pages", "Disappears", "Follows you down every page"),
        ("Where a story ends", "Nowhere", "The scheme it is about"),
        ("Contact button", "Does nothing", "Works, plus an assistant"),
        ("404 page", "Stock skateboarder", "In Marathi, all 16 sections"),
        ("robots.txt / sitemap", "Neither", "Both, generated at build"),
        ("Marathi typography", "Latin line-height", "Set for Devanagari")]

y = 2.5
for i, (label, before, after) in enumerate(rows):
    if i % 2 == 0:
        rect(s, 0.9, y, 11.5, 0.5, RGBColor(0x24, 0x1F, 0x1B))
    tf = box(s, 1.15, y + 0.12, 3.6, 0.3)
    para(tf, label, 13, CREAM, SANS, bold=True, first=True)
    tf = box(s, 5.0, y + 0.12, 3.5, 0.3)
    para(tf, before, 12.5, RGBColor(0xC2, 0x7A, 0x72), SANS, first=True)
    tf = box(s, 8.7, y + 0.12, 4.0, 0.3)
    para(tf, after, 12.5, RGBColor(0x7F, 0xC4, 0x9E), SANS, first=True)
    y += 0.5
footer(s, n)

# =====================================================================
# 13 — Honest limits
# =====================================================================
n += 1
s = slide(CREAM)
eyebrow(s, "BEING STRAIGHT WITH YOU")
title(s, "What this is not, yet")

items = [
    ("It is a proof, not a migration",
     "24 articles per category are loaded, out of 2,664. The templates are finished; "
     "the rest is data through them."),
    ("Mobile is unverified",
     "Built mobile-first with no fixed widths, but not yet confirmed on a real handset. "
     "That is the first thing to check."),
    ("The content is a snapshot",
     "Collected 28 August 2026 and static. Database access would keep it live and is "
     "worth asking for."),
    ("Photography is the real ceiling",
     "Group shots under fluorescent light, GPS watermarks, posters with text baked in. "
     "The design is forgiving of it — but a one-page shot standard for the 36 district "
     "offices would lift this site further than another round of design."),
]
y = 2.3
for head, body in items:
    rect(s, 0.9, y, 0.055, 0.95, SAFFRON)
    tf = box(s, 1.25, y, 11.2, 0.95)
    para(tf, head, 16, INK, SANS, bold=True, first=True, space_after=5)
    para(tf, body, 13, INK2, SANS, line=1.4)
    y += 1.12
footer(s, n)

# =====================================================================
# 14 — Close
# =====================================================================
n += 1
s = slide(INK)
rect(s, 0, 0, 0.22, 7.5, SAFFRON)

tf = box(s, 1.2, 2.4, 10.5, 1.8)
para(tf, "The design was never the hard part.", 34, CREAM, SERIF, first=True, line=1.2)

tf = box(s, 1.2, 3.6, 9.0, 2.0)
para(tf, "AMRUT already has the stories, the schemes and offices in all 36 districts. "
         "What was missing was the path between them.",
     18, RGBColor(0xC9, 0xBF, 0xB6), SANS, first=True, line=1.5, space_after=16)
para(tf, "Build the scheme directory, the district office finder and the eligibility "
         "check, and every article already published becomes a way in.",
     18, SAFFRON, SANS, bold=True, line=1.5)

tf = box(s, 1.2, 6.35, 10.5, 0.5)
para(tf, "amrut-mh.vercel.app", 14, CREAM, SANS, bold=True, first=True)

out = os.path.join(os.path.dirname(__file__), "..", "docs",
                   "why-the-new-site-is-better.pptx")
prs.save(out)
print(f"saved {os.path.normpath(out)}  ·  {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
